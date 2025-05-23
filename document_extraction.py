from io import BufferedReader
import PyPDF2
import docx2txt
import json
import xmltodict
from PIL import Image
import pytesseract
import textract
import subprocess
from ebooklib import epub
from ebooklib.epub import EpubHtml
from lxml import html
from pptx import Presentation
# import pandas as pd



class TextExtractor:

    ALL_POSSIBLE_EXTENSIONS = ["pdf", "docx", "doc", "txt", "json", "xml", "jpg", "png",
                                "jpeg", "tiff", "djvu", "epub", "pptx", "xlsx", "xls",
]

    def __init__(self, document_path: str = None) -> None:

        try:
            if document_path and self._verify_if_eligible(document_path):
                self.document_path = document_path
                self.extension = document_path.split(".")[-1]
            else:
                self.document_path = None
                self.extension = None

        except ValueError as e:
            print(e)

    def assign_new_document(self, document_path: str) -> bool:
        try:
            if self._verify_if_eligible(document_path):
                self.document_path = document_path
                self.extension = document_path.split(".")[-1]
                return True
        except ValueError as e:
            raise e

    def _verify_if_eligible(self, document_path: str) -> bool:

        if document_path is None:
            raise DocumentNotSupported("Document path cannot be None")
        if not isinstance(document_path, str):
            raise DocumentNotSupported("Document path must be a string")
        if not document_path.split(".")[-1] in self.ALL_POSSIBLE_EXTENSIONS:
            raise DocumentNotSupported("Document format not supported")

        return True

    def extract_text(self, get_dict: bool = False, get_image_blocks: bool = False, warning: bool = False) -> str:
        if self.document_path is None:
            return "No document path assigned"

        if self.extension in ["pdf", "docx", "doc", "txt"] and get_dict and warning:
            print(
                "Warning: get_dict is only supported for json and xml files, ignoring get_dict")

        if self.extension == "pdf":
            return self._extract_text_from_pdf()
        elif self.extension == "djvu":
            return self._extract_text_from_djvu()
        elif self.extension == "pptx":
            return self._extract_text_from_pptx()
        elif self.extension in ["xlsx", "xls"]:
            return self._extract_text_from_xlsx()
        elif self.extension == "epub":
            return self._extract_text_from_epub()
        elif self.extension == "doc":
            return self._extract_text_from_doc()
        elif self.extension == "docx":
            return self._extract_text_from_docx()
        elif self.extension == "txt":
            return self._extract_text_from_txt()
        elif self.extension == "json":
            return self._extract_text_from_json(get_dict)
        elif self.extension == "xml":
            return self._extract_text_from_xml(get_dict)
        elif self.extension in ["jpg", "png", "jpeg", "tiff"]:
            return self._extract_text_from_picture(get_image_blocks)
        else:
            return "Extension not supported"

    def _extract_text_from_pdf(self) -> str:
        file_path: str = self.document_path
        pdf_file_obj: BufferedReader = open(file_path, 'rb')
        pdf_reader: PyPDF2.PdfReader = PyPDF2.PdfReader(pdf_file_obj)
        text: str = ''
        for page in pdf_reader.pages:
            text += page.extract_text()
        pdf_file_obj.close()
        return text

    def _extract_text_from_djvu(self):
        file_path: str = self.document_path
        output_djvu = subprocess.run(['djvutxt', file_path], capture_output=True, text=True)
        text = output_djvu.stdout
        return text

    def _extract_text_from_epub(self):
        file_path: str = self.document_path
        book = epub.read_epub(file_path)
        text = ''
        for item in book.get_items():
            if isinstance(item, EpubHtml):
                tree = html.fromstring(item.get_body_content())
                text += tree.text_content()
        return text

    def _extract_text_from_doc(self):
        file_path: str = self.document_path
        text = textract.process(file_path).decode("utf-8")
        return text

    def _extract_text_from_docx(self) -> str:
        file_path: str = self.document_path
        text: str = docx2txt.process(file_path)
        return text

    def _extract_text_from_txt(self) -> str:
        file_path: str = self.document_path
        with open(file_path, 'r') as file:
            text: str = file.read()
        return text

    def _extract_text_from_json(self, get_dict: bool) -> str:
        file_path: str = self.document_path
        with open(file_path, 'r') as file:
            if get_dict:
                return json.load(file)
            else:
                return file.read()

    def _extract_text_from_xml(self, get_dict: bool) -> str:
        file_path: str = self.document_path
        with open(file_path, 'r') as file:
            if get_dict:
                return xmltodict.parse(file.read())
            else:
                return file.read()

    def _tesseract_installed(self) -> bool:
        try:
            pytesseract.get_tesseract_version()
            return True
        except pytesseract.TesseractNotFoundError:
            return False

    def _extract_text_from_picture(self, get_image_blocks: bool) -> str:    
        file_path: str = self.document_path
        image = Image.open(file_path)
        if not self._tesseract_installed():
            raise TesseractNotInstalled("Tesseract not installed")
        if not get_image_blocks:
            text: str = pytesseract.image_to_string(image, lang="rus+eng")
        if get_image_blocks:
            text: str = pytesseract.image_to_data(image, lang="rus+eng", output_type=pytesseract.Output.DICT)
        return text
    # def _extract_text_from_xlsx(self) -> str:
    #     file_path: str = self.document_path
    #     text = ""
    #     # Чтение всех листов
    #     # with pd.ExcelFile(file_path) as xls:
    #     #     for sheet_name in xls.sheet_names:
    #     #         df = pd.read_excel(xls, sheet_name=sheet_name)
    #     #         print(f"Sheet: {sheet_name}")
    #     #         print(df.to_string())  # Вывод текста из таблицы
    #     #         text += df.to_string()
    #     #         text += '\n\n'

    #     return text


    def _extract_text_from_pptx(self) -> str:
        file_path: str = self.document_path
        presentation = Presentation(file_path)
        extracted_text = ""

        for slide_number, slide in enumerate(presentation.slides):
            extracted_text += f"\nSlide {slide_number + 1}:\n"

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + "\n"

        return extracted_text



class DocumentNotSupported(Exception):
    def __init__(self, message="Document format not supported"):
        self.message = message
        super().__init__(self.message)


class TesseractNotInstalled(Exception):
    def __init__(self, message="Tesseract not installed"):
        self.message = message
        super().__init__(self.message)

